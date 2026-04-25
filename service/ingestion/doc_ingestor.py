"""
service/ingestion/doc_ingestor.py
===================================
Ingests DOCX and PPTX files into RawDocuments.

DOCX: Extract paragraph text using python-docx
PPTX: Extract slide text using python-pptx, one RawDocument per slide

Used by:
    Onboarding pipeline for document file uploads
"""

import base64
import io

from langsmith import traceable

from core.logging import get_logger, timer
from llm.llm_schema.state_models import Modality, SourceType
from service.ingestion.chunker import RawDocument

logger = get_logger(__name__)


class DocIngestor:
    """
    Extracts text from DOCX and PPTX files.

    DOCX: All paragraphs joined into a single RawDocument
          (chunker handles splitting into token-sized chunks)

    PPTX: Each slide becomes one RawDocument with slide number
          preserved in metadata
    """

    @traceable(name="doc_ingestor", run_type="retriever")
    async def ingest(
        self,
        files: list[dict]  # [{"filename": str, "base64_content": str, "mime_type": str}]
    ) -> list[RawDocument]:
        """
        Ingest a list of base64-encoded DOCX or PPTX files.

        Routes each file to the appropriate extractor based
        on filename extension or MIME type.

        Args:
            files: list of FileUpload dicts

        Returns:
            list of RawDocument objects
        """
        documents: list[RawDocument] = []

        for file in files:
            filename = file.get("filename", "unknown_doc")
            mime_type = file.get("mime_type", "")

            with timer("doc_extract", logger, extra={"filename": filename}):
                try:
                    raw_bytes = base64.b64decode(file["base64_content"])

                    if self._is_docx(filename, mime_type):
                        docs = await self._ingest_docx(raw_bytes, filename)
                    elif self._is_pptx(filename, mime_type):
                        docs = await self._ingest_pptx(raw_bytes, filename)
                    else:
                        logger.warning(
                            "Unsupported document format — skipping",
                            extra={"filename": filename, "mime_type": mime_type}
                        )
                        continue

                    documents.extend(docs)

                except Exception as e:
                    logger.warning(
                        "Document ingestion failed — skipping",
                        extra={"filename": filename, "error": str(e)}
                    )

        logger.info(
            "Document ingestion complete",
            extra={
                "files_processed": len(files),
                "documents_extracted": len(documents),
            }
        )
        return documents

    async def _ingest_docx(
        self,
        file_bytes: bytes,
        origin: str
    ) -> list[RawDocument]:
        """
        Extract text from a DOCX file.
        Joins all non-empty paragraphs into a single text block.
        """
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))

        # Extract title from core properties
        title = None
        try:
            title = doc.core_properties.title or None
        except Exception:
            pass

        # Join all non-empty paragraphs
        paragraphs = [
            p.text.strip()
            for p in doc.paragraphs
            if p.text.strip()
        ]

        if not paragraphs:
            logger.warning(
                "DOCX has no extractable text",
                extra={"origin": origin}
            )
            return []

        full_text = "\n\n".join(paragraphs)

        return [RawDocument(
            source_type=SourceType.DOCX,
            origin=origin,
            title=title,
            text=full_text,
            modality=Modality.TEXT,
        )]

    async def _ingest_pptx(
        self,
        file_bytes: bytes,
        origin: str
    ) -> list[RawDocument]:
        """
        Extract text from a PPTX file.
        Each slide becomes one RawDocument with slide number.
        """
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        documents: list[RawDocument] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            # Extract all text from slide shapes
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())

            if not texts:
                continue  # Skip blank slides

            slide_text = "\n".join(texts)

            documents.append(RawDocument(
                source_type=SourceType.PPTX,
                origin=origin,
                title=f"Slide {slide_num}",
                text=slide_text,
                modality=Modality.TEXT,
                slide=slide_num,
            ))

        if not documents:
            logger.warning(
                "PPTX has no extractable text",
                extra={"origin": origin}
            )

        return documents

    @staticmethod
    def _is_docx(filename: str, mime_type: str) -> bool:
        return (
            filename.lower().endswith(".docx") or
            "wordprocessingml" in mime_type or
            mime_type == "application/vnd.openxmlformats-officedocument"
                         ".wordprocessingml.document"
        )

    @staticmethod
    def _is_pptx(filename: str, mime_type: str) -> bool:
        return (
            filename.lower().endswith(".pptx") or
            "presentationml" in mime_type or
            mime_type == "application/vnd.openxmlformats-officedocument"
                         ".presentationml.presentation"
        )
    